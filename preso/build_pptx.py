#!/usr/bin/env python3
"""Build a .pptx deck deterministically from a Markdown speaker-script.

The Markdown is the source of record (see preso/README.md); this script only *mirrors*
it into slides. It is intentionally small and dependency-light: the only third-party
import is python-pptx, and it is guarded so a missing dep produces an actionable message
rather than a traceback.

Deck format (deliberately minimal and predictable):

    # Title of the deck                 <- first slide: title (+ optional '> subtitle')
    > one-line subtitle

    ---                                 <- a line that is exactly '---' starts a new slide

    ## Slide heading
    > BLUF line shown under the heading
    - a bullet
    - another bullet

    Notes:                              <- everything after a 'Notes:' line = speaker notes
    The precise, caveated version, and the okf/ note(s) this slide traces back to.

Usage:
    python3 build_pptx.py <deck.md> [--out <deck.pptx>] [--template <branded.pptx>]

Exit codes: 0 ok; 2 usage/IO error; 3 python-pptx not installed.
"""
from __future__ import annotations

import argparse
import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ModuleNotFoundError:
    sys.stderr.write(
        "error: python-pptx is not installed.\n"
        "       Deck building is the only part of this repo that needs a third-party\n"
        "       dependency. Install it with:\n\n"
        "           pip install -r preso/requirements.txt\n\n"
    )
    raise SystemExit(3)


def parse_deck(text: str):
    """Return (title, subtitle, title_notes, [slides]) where each slide is a dict with
    heading, bluf, bullets, body, notes. Slides are separated by lines that are exactly
    '---'."""
    blocks = []
    current: list[str] = []
    for line in text.splitlines():
        if line.strip() == "---":
            blocks.append(current)
            current = []
        else:
            current.append(line)
    blocks.append(current)

    slides = []
    deck_title, deck_subtitle, deck_notes = "Untitled deck", "", ""

    for i, block in enumerate(blocks):
        heading, bluf, bullets, body, notes = "", "", [], [], []
        in_notes = False
        for raw in block:
            line = raw.rstrip()
            if not line.strip():
                continue
            if line.strip() == "Notes:":
                in_notes = True
                continue
            if in_notes:
                notes.append(line.strip())
                continue
            if line.startswith("## "):
                heading = line[3:].strip()
            elif line.startswith("# "):
                heading = line[2:].strip()
            elif line.startswith("> "):
                bluf = (bluf + " " + line[2:].strip()).strip()
            elif line.startswith(("- ", "* ")):
                bullets.append(line[2:].strip())
            else:
                body.append(line.strip())

        if i == 0:
            # first block is the title slide
            deck_title = heading or deck_title
            deck_subtitle = bluf or " ".join(body)
            deck_notes = "\n".join(notes)
            continue
        if not (heading or bullets or body or notes):
            continue  # skip empty trailing block
        slides.append(
            {"heading": heading, "bluf": bluf, "bullets": bullets,
             "body": body, "notes": "\n".join(notes)}
        )

    return deck_title, deck_subtitle, deck_notes, slides


def build(md_path: str, out_path: str, template: str | None) -> int:
    with open(md_path, encoding="utf-8") as f:
        title, subtitle, title_notes, slides = parse_deck(f.read())

    prs = Presentation(template) if template else Presentation()

    # Title slide
    title_layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(title_layout)
    s.shapes.title.text = title
    if len(s.placeholders) > 1 and subtitle:
        s.placeholders[1].text = subtitle
    if title_notes:
        s.notes_slide.notes_text_frame.text = title_notes

    # Content slides: use a blank-ish layout and draw a title + a bullet body box.
    content_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    for slide in slides:
        sl = prs.slides.add_slide(content_layout)
        if sl.shapes.title is not None:
            sl.shapes.title.text = slide["heading"]

        lines = []
        if slide["bluf"]:
            lines.append((slide["bluf"], True))
        lines.extend((b, False) for b in slide["bullets"])
        lines.extend((b, False) for b in slide["body"])

        body_ph = None
        for ph in sl.placeholders:
            if ph.placeholder_format.idx == 1:
                body_ph = ph
                break
        if body_ph is None:
            body_ph = sl.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(8.6), Inches(4.8))

        tf = body_ph.text_frame
        tf.clear()
        if lines:
            first_text, first_bold = lines[0]
            tf.paragraphs[0].text = first_text
            tf.paragraphs[0].font.bold = first_bold
            for text, bold in lines[1:]:
                p = tf.add_paragraph()
                p.text = text
                p.font.bold = bold
                p.font.size = Pt(18)

        if slide["notes"]:
            sl.notes_slide.notes_text_frame.text = slide["notes"]

    prs.save(out_path)
    print(f"built {out_path}  ({len(slides)} content slide(s) + title)")
    return 0


def main() -> int:
    ap = argparse.ArgumentParser(description="Build a .pptx deck from a Markdown speaker-script.")
    ap.add_argument("deck", help="path to the deck .md")
    ap.add_argument("--out", help="output .pptx (default: alongside the .md)")
    ap.add_argument("--template", help="branded .pptx/.potx to build on top of")
    args = ap.parse_args()

    if not os.path.isfile(args.deck):
        sys.stderr.write(f"error: no such deck file: {args.deck}\n")
        return 2
    out = args.out or os.path.splitext(args.deck)[0] + ".pptx"
    template = args.template
    if template and not os.path.isfile(template):
        sys.stderr.write(f"error: template not found: {template}\n")
        return 2
    return build(args.deck, out, template)


if __name__ == "__main__":
    raise SystemExit(main())
